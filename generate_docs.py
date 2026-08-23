import sys
import subprocess
import os

# Install required packages
def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    import docx
except ImportError:
    print("Installing python-docx...")
    install("python-docx")

try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    install("python-pptx")

from docx import Document
from docx.shared import Inches, Pt
from docx.shared import RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PtInches, Pt as PtFont
from pptx.dml.color import RGBColor as PptxRGBColor

def create_word_report(file_path):
    doc = Document()
    
    # Configure styles
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run("PROJECT REPORT\nINFI AI VOICE ASSISTANT")
    run.font.name = 'Arial'
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = DocxRGBColor(99, 102, 241) # Indigo color matching UI

    subtitle_p = doc.add_paragraph()
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_sub = subtitle_p.add_run("A Next-Generation Multiplatform Voice AI Workspace")
    run_sub.font.name = 'Arial'
    run_sub.font.size = Pt(14)
    run_sub.font.italic = True
    
    doc.add_paragraph("\n" * 2)
    
    # Section: What
    doc.add_heading("1. What is INFI Voice Assistant?", level=1)
    doc.add_paragraph(
        "INFI is an advanced, high-fidelity AI-powered voice assistant capable of running across desktop and mobile devices. "
        "It features a premium, responsive glassmorphic dark-theme UI designed to engage users with micro-animations and real-time audio visualization. "
        "By utilizing the Google Gemini API for cognitive tasks and a FastAPI backend for local desktop execution, "
        "INFI bridges the gap between pure digital conversation and operating system-level control."
    )
    
    # Section: Why
    doc.add_heading("2. Why we built INFI?", level=1)
    doc.add_paragraph(
        "Modern virtual assistants are either locked into proprietary operating systems (like Siri or Cortana) or operate entirely "
        "as browser text-interfaces without desktop automation capability. INFI was created to achieve three primary goals:\n"
        "• Hands-free productivity: Enabling users to launch local applications, websites, control volume, and take screenshots using voice.\n"
        "• Unified Multiplatform Core: Operating fully in the browser with local system fallbacks so that if the desktop is disconnected, the assistant remains useful on mobile.\n"
        "• Accurate Problem Solving: Using advanced large language models (LLMs) like Gemini-3.5-flash for complex cognitive queries while executing mathematical equations locally to eliminate AI hallucinations."
    )
    
    # Section: How
    doc.add_heading("3. How does INFI work? (System Architecture)", level=1)
    doc.add_paragraph(
        "INFI is structured as a two-tier fullstack application:\n"
        "1. Frontend (React + Vite): Handles speech recognition (via browser Web Speech API), local calculations, state management, "
        "and settings (Gemini API credentials). It is styled with vanilla CSS leveraging modern dark themes and glowing animations.\n"
        "2. Backend (FastAPI + Python): A local service running on the user's desktop. It listens on port 5000 and executes system-level operations "
        "such as launching applications (VS Code, Chrome, etc.), opening URLs, locking the desktop workstation, adjusting master volume, "
        "and reading hardware performance metrics (CPU and memory usage).\n"
        "3. Brain (Google Gemini API): When general queries, programming requests, or complex logic are initiated, the frontend calls the Gemini API directly using the configured credentials."
    )
    
    # Section: Key Features
    doc.add_heading("4. Key Features Implemented", level=1)
    features = [
        ("🎙️ Real-time Voice Control", "Implements Web Speech API with TTS voice selection and real-time canvas visualizer."),
        ("🎵 Platform Music Streaming", "Intelligently routes song play requests to Spotify (app/web) or YouTube depending on user commands (e.g., 'play song in spotify')."),
        ("🖥️ Native OS Automation", "Adjusts volume, takes screenshots, locks screen, and launches applications via python subprocesses."),
        ("💻 Built-in Code Playground", "Generates code snippets in Python, JS, C++, Go, and CSS with syntax highlighting and quick-copy option."),
        ("📈 Hardware Performance monitor", "Displays real-time CPU and RAM health metrics directly within the assistant sidebar.")
    ]
    
    for title, desc in features:
        p = doc.add_paragraph(style='List Bullet')
        r_t = p.add_run(f"{title}: ")
        r_t.bold = True
        p.add_run(desc)
        
    doc.save(file_path)
    print(f"Word report created successfully at {file_path}")
 
def create_ppt_presentation(file_path):
    prs = Presentation()
    
    c_indigo = PptxRGBColor(99, 102, 241)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "INFI Voice Assistant"
    subtitle.text = "A Next-Generation Multiplatform Voice AI Workspace\nBuilt with React, FastAPI & Google Gemini API"
    title.text_frame.paragraphs[0].font.color.rgb = c_indigo
    
    # Slide 2: What is INFI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "What is INFI?"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "INFI is a hybrid desktop-web AI Voice Assistant."
    p1 = tf.add_paragraph()
    p1.text = "• Premium glassmorphic interface with canvas-based audio visualizer."
    p2 = tf.add_paragraph()
    p2.text = "• Utilizes Gemini API for logic, coding, and general intelligence."
    p3 = tf.add_paragraph()
    p3.text = "• Connects to a local Python daemon for native Windows control."
    
    # Slide 3: Why we built it
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Why we built INFI?"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Bridge the gap between OS control and AI reasoning:"
    p1 = tf.add_paragraph()
    p1.text = "• Speed up tasks: hands-free opening of apps, websites, and volume adjustment."
    p2 = tf.add_paragraph()
    p2.text = "• Accuracy: Local execution of mathematical inputs preventing AI hallucination."
    p3 = tf.add_paragraph()
    p3.text = "• Multiplatform: Access via browser anywhere with fallback web actions on mobile."
    
    # Slide 4: System Architecture (How)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture (How it works)"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Frontend + Backend + Cloud LLM integration:"
    p1 = tf.add_paragraph()
    p1.text = "• React Frontend: Manages TTS/STT, displays logs, and handles custom regex router."
    p2 = tf.add_paragraph()
    p2.text = "• FastAPI Backend: Listens on port 5000, runs OS subprocesses, takes screenshots."
    p3 = tf.add_paragraph()
    p3.text = "• Google Gemini API: Performs coding, analysis, and answers questions on the fly."
    
    # Slide 5: Core Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Core Features"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Key capabilities of the INFI workspace:"
    p1 = tf.add_paragraph()
    p1.text = "• 🎙️ Voice & Audio controls with visual neon waveform feedback."
    p2 = tf.add_paragraph()
    p2.text = "• 🎵 Platform-smart music play (YouTube vs. Spotify)."
    p3 = tf.add_paragraph()
    p3.text = "• 💻 Interactive Coding Solver Playground with code-copy options."
    p4 = tf.add_paragraph()
    p4.text = "• 🖥️ System actions: Adjust volume, mute, take screenshots, lock workstation."
    
    # Slide 6: Summary & Deployment
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deployment & Live Status"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Active links and project workspace:"
    p1 = tf.add_paragraph()
    p1.text = "• Codebase hosted on GitHub: IfrahAnsari/my_ai_assistant"
    p2 = tf.add_paragraph()
    p2.text = "• Deployed Live Demo URL: https://myaiassistant-tau.vercel.app"
    p3 = tf.add_paragraph()
    p3.text = "• Prepopulated settings and API key details for out-of-the-box performance."

    prs.save(file_path)
    print(f"PowerPoint created successfully at {file_path}")

if __name__ == "__main__":
    desktop_dir = r"c:\Users\FAREA\OneDrive\Desktop\my_ai_assistant"
    create_word_report(os.path.join(desktop_dir, "INFI_Voice_Assistant_Report.docx"))
    create_ppt_presentation(os.path.join(desktop_dir, "INFI_Voice_Assistant_Presentation.pptx"))
